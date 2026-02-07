"""
Office Agent
============
PPT, Word, Excel 문서 자동 생성 에이전트

참고: Reference/Gemini-Claw/src/office/office.py
"""

import os
import logging
from typing import List, Dict, Optional
from datetime import datetime

# Word
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn

# PowerPoint
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor

# Excel
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

logger = logging.getLogger(__name__)


class OfficeAgent:
    """
    Office 문서 자동 생성 에이전트
    
    기능:
    - Word 보고서 생성 (구조화된 섹션)
    - PowerPoint 프레젠테이션 생성
    - Excel 스프레드시트 생성
    - 한국어 폰트 지원
    """
    
    # 공통 스타일 설정
    STYLE_CONFIG = {
        "font_kr": "Malgun Gothic",
        "font_en": "Calibri",
        "primary_color": "366092",
        "accent_color": "4472C4",
        "header_bg": "366092",
        "header_fg": "FFFFFF",
    }
    
    def __init__(self, output_dir: str = "output"):
        """
        Args:
            output_dir: 출력 파일 저장 디렉토리
        """
        self.output_dir = os.path.abspath(output_dir)
        self.default_font = self.STYLE_CONFIG["font_kr"]
        
        # 출력 디렉토리 생성
        os.makedirs(self.output_dir, exist_ok=True)
        logger.info(f"[OfficeAgent] Initialized. Output dir: {self.output_dir}")
    
    def _get_output_path(self, filename: str) -> str:
        """출력 파일 경로 생성"""
        return os.path.join(self.output_dir, filename)
    
    def _apply_korean_font_word(self, run):
        """Word 문서에 한국어 폰트 적용"""
        run.font.name = self.default_font
        # East Asian 폰트 설정 (한글 지원)
        run._element.rPr.rFonts.set(qn('w:eastAsia'), self.default_font)
    
    # =========================================================================
    # PowerPoint 생성
    # =========================================================================
    
    def create_presentation(
        self,
        title: str,
        subtitle: str = "",
        slides: Optional[List[Dict]] = None,
        output_path: str = "presentation.pptx"
    ) -> Dict:
        """
        PowerPoint 프레젠테이션 생성
        
        Args:
            title: 발표 제목
            subtitle: 부제목
            slides: 슬라이드 데이터 목록
                    [{"title": "슬라이드 제목", "content": ["항목1", "항목2"]}]
            output_path: 출력 파일명
        
        Returns:
            {"success": bool, "path": str, "message": str}
        """
        try:
            prs = Presentation()
            slides = slides or []
            
            # 1. 타이틀 슬라이드 (Layout 0)
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            if subtitle and len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = subtitle
            
            # 폰트 적용
            self._apply_font_to_slide(title_slide)
            
            # 2. 콘텐츠 슬라이드 (Layout 1 - Title and Content)
            for slide_data in slides:
                slide_title = slide_data.get("title", "Untitled")
                slide_content = slide_data.get("content", [])
                
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_title
                
                # 콘텐츠 영역
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    
                    if isinstance(slide_content, list):
                        for i, item in enumerate(slide_content):
                            if i == 0:
                                tf.text = str(item)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(item)
                                p.level = 0
                    elif isinstance(slide_content, str):
                        tf.text = slide_content
                
                self._apply_font_to_slide(slide)
            
            # 저장
            full_path = self._get_output_path(output_path)
            prs.save(full_path)
            
            logger.info(f"[OfficeAgent] PPT created: {full_path}")
            return {
                "success": True,
                "path": full_path,
                "message": f"Successfully created presentation with {len(slides) + 1} slides",
                "slide_count": len(slides) + 1
            }
            
        except Exception as e:
            logger.error(f"[OfficeAgent] PPT creation failed: {e}")
            return {
                "success": False,
                "path": "",
                "message": f"Error creating presentation: {str(e)}"
            }
    
    def _apply_font_to_slide(self, slide):
        """슬라이드 전체에 폰트 적용"""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = self.default_font
                for run in paragraph.runs:
                    run.font.name = self.default_font
    
    # =========================================================================
    # Word 문서 생성
    # =========================================================================
    
    def create_word_report(
        self,
        title: str,
        sections: Optional[List[Dict[str, str]]] = None,
        output_path: str = "report.docx"
    ) -> Dict:
        """
        Word 보고서 생성
        
        Args:
            title: 문서 제목
            sections: 섹션 목록
                      [{"heading": "섹션 제목", "content": "본문 내용"}]
            output_path: 출력 파일명
        
        Returns:
            {"success": bool, "path": str, "message": str}
        """
        try:
            doc = Document()
            sections = sections or []
            
            # 제목
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in title_para.runs:
                self._apply_korean_font_word(run)
            
            # 섹션들
            for section in sections:
                heading_text = section.get("heading", "")
                content_text = section.get("content", "")
                
                # 섹션 헤딩
                if heading_text:
                    h = doc.add_heading(heading_text, level=1)
                    for run in h.runs:
                        self._apply_korean_font_word(run)
                
                # 본문 내용
                if content_text:
                    for para_text in content_text.split('\n'):
                        if para_text.strip():
                            p = doc.add_paragraph(para_text.strip())
                            for run in p.runs:
                                self._apply_korean_font_word(run)
            
            # 저장
            full_path = self._get_output_path(output_path)
            doc.save(full_path)
            
            logger.info(f"[OfficeAgent] Word created: {full_path}")
            return {
                "success": True,
                "path": full_path,
                "message": f"Successfully created Word report with {len(sections)} sections",
                "section_count": len(sections)
            }
            
        except Exception as e:
            logger.error(f"[OfficeAgent] Word creation failed: {e}")
            return {
                "success": False,
                "path": "",
                "message": f"Error creating Word report: {str(e)}"
            }
    
    # =========================================================================
    # Excel 생성
    # =========================================================================
    
    def create_excel(
        self,
        data: List[Dict],
        output_path: str = "data.xlsx",
        sheet_name: str = "Data"
    ) -> Dict:
        """
        Excel 스프레드시트 생성
        
        Args:
            data: 데이터 목록 (각 dict가 한 행)
                  [{"Name": "Kim", "Age": 30}, {"Name": "Lee", "Age": 25}]
            output_path: 출력 파일명
            sheet_name: 시트 이름
        
        Returns:
            {"success": bool, "path": str, "message": str}
        """
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            if not data:
                return {
                    "success": False,
                    "path": "",
                    "message": "No data provided"
                }
            
            # 헤더 (첫 번째 dict의 키들)
            headers = list(data[0].keys())
            
            # 스타일 정의
            header_font = Font(name=self.default_font, bold=True, color=self.STYLE_CONFIG["header_fg"])
            header_fill = PatternFill(start_color=self.STYLE_CONFIG["header_bg"], 
                                       end_color=self.STYLE_CONFIG["header_bg"], 
                                       fill_type="solid")
            body_font = Font(name=self.default_font)
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # 헤더 작성
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border
            
            # 데이터 작성
            for row_idx, row_data in enumerate(data, 2):
                for col_idx, header in enumerate(headers, 1):
                    value = row_data.get(header, "")
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.font = body_font
                    cell.border = thin_border
            
            # 열 너비 자동 조정
            for col_idx, header in enumerate(headers, 1):
                max_length = len(str(header))
                for row_data in data:
                    cell_value = str(row_data.get(header, ""))
                    max_length = max(max_length, len(cell_value))
                ws.column_dimensions[get_column_letter(col_idx)].width = min(max_length + 2, 50)
            
            # 저장
            full_path = self._get_output_path(output_path)
            wb.save(full_path)
            
            logger.info(f"[OfficeAgent] Excel created: {full_path}")
            return {
                "success": True,
                "path": full_path,
                "message": f"Successfully created Excel with {len(data)} rows",
                "row_count": len(data)
            }
            
        except Exception as e:
            logger.error(f"[OfficeAgent] Excel creation failed: {e}")
            return {
                "success": False,
                "path": "",
                "message": f"Error creating Excel: {str(e)}"
            }


# 간단한 테스트
if __name__ == "__main__":
    agent = OfficeAgent(output_dir="output")
    
    # PPT 테스트
    result = agent.create_presentation(
        title="Tiny-MoA 소개",
        subtitle="GPU Poor를 위한 AI 군단",
        slides=[
            {"title": "개요", "content": ["멀티 에이전트 시스템", "1.2B Thinking Model", "경량화된 설계"]},
            {"title": "핵심 기능", "content": ["Tool Calling", "RAG", "Office 문서 생성"]},
        ],
        output_path="test_presentation.pptx"
    )
    print(result)
